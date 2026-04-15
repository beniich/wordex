from fastapi import APIRouter, Depends, HTTPException, Body
from fastapi.responses import StreamingResponse
import asyncpg
from typing import Dict, Any
import json
import uuid
import io

from app.database import get_db
from app.auth import get_current_user_id
from app.routers.documents import update_document, get_document, list_versions, restore_version
from app.models import DocumentUpdate
from app.routers.ai import ollama_generate, SYSTEM_PROMPTS

router = APIRouter()

@router.get("/{id}")
async def get_presentation(
    id: str,
    user_id: str = Depends(get_current_user_id),
    db: asyncpg.Connection = Depends(get_db)
):
    """Retrieve a presentation document."""
    doc = await get_document(id, user_id=user_id, db=db)
    if doc["doc_type"] != "presentation":
        raise HTTPException(400, "Document is not a presentation")
    
    content = doc["content"] or {
        "title": doc["title"],
        "slides": [],
        "theme": "corporate",
        "estimated_duration": 0
    }
    return content

@router.put("/{id}")
async def update_presentation(
    id: str,
    body: Dict[str, Any],
    user_id: str = Depends(get_current_user_id),
    db: asyncpg.Connection = Depends(get_db)
):
    """Update a presentation slides and metadata."""
    update_data = DocumentUpdate(content=body, title=body.get("title", None))
    updated_doc = await update_document(id, update_data, user_id=user_id, db=db)
    return {"status": "updated", "version": await db.fetchval("SELECT count(*) FROM document_versions WHERE document_id=$1", id)}

@router.get("/{id}/versions")
async def get_presentation_versions(
    id: str,
    user_id: str = Depends(get_current_user_id),
    db: asyncpg.Connection = Depends(get_db)
):
    return await list_versions(id, user_id=user_id, db=db)

@router.post("/{id}/restore/{version_id}")
async def restore_presentation_version(
    id: str,
    version_id: str,
    user_id: str = Depends(get_current_user_id),
    db: asyncpg.Connection = Depends(get_db)
):
    return await restore_version(id, version_id, user_id=user_id, db=db)

@router.post("/{id}/generate-from-ai")
async def generate_slides_from_ai(
    id: str,
    body: dict = Body(...),
    user_id: str = Depends(get_current_user_id),
    db: asyncpg.Connection = Depends(get_db)
):
    """Generate slides full content from AI and save to doc."""
    doc = await get_document(id, user_id=user_id, db=db)
    if doc["doc_type"] != "presentation":
        raise HTTPException(400, "Document is not a presentation")

    topic = body.get("topic", "")
    n_slides = body.get("n_slides", 8)

    prompt = f'''
# TASK: Create a high-quality professional presentation
Topic: {topic}
Number of slides: {n_slides}
Theme: sable

# CONSTRAINTS:
1. Return ONLY pure valid JSON. 
2. NO markdown blocks (```json ... ```).
3. NO conversational filler.
4. Language: Match the topic/input language.

# STRUCTURE:
{{
  "title": "Main Title",
  "theme": "sable",
  "slides": [
    {{
      "id": "s1",
      "title": "Slide Title",
      "content": "Professional content...",
      "speakerNotes": "Context for the presenter",
      "layout": "titleAndContent",
      "visualType": "bulletPoints"
    }}
  ]
}}
'''

    result = await ollama_generate(
        prompt=prompt,
        system=SYSTEM_PROMPTS.get("designer", ""),
        temperature=0.2,
    )
    
    # Try parsing the result JSON
    try:
        if "```json" in result:
            json_str = result.split("```json")[1].split("```")[0].strip()
        elif "```" in result:
            json_str = result.split("```")[1].split("```")[0].strip()
        else:
            json_str = result.strip()
            
        presentation_data = json.loads(json_str)
        
        for slide in presentation_data.get("slides", []):
            if "id" not in slide or slide["id"] == "s1" or slide["id"] == "s2":
                slide["id"] = str(uuid.uuid4())
                
    except Exception as e:
        raise HTTPException(500, f"AI returned invalid JSON: {str(e)} -> Content: {result[:100]}...")

    update_data = DocumentUpdate(
        content=presentation_data, 
        title=presentation_data.get("title", doc["title"])
    )
    await update_document(id, update_data, user_id=user_id, db=db)
    
    return {"status": "success", "presentation": presentation_data}

@router.post("/{id}/generate-from-doc")
async def generate_slides_from_doc(
    id: str,
    body: dict = Body(...),
    user_id: str = Depends(get_current_user_id),
    db: asyncpg.Connection = Depends(get_db)
):
    """Transform an existing document (note) into a presentation."""
    source_id = body.get("source_doc_id")
    if not source_id:
        raise HTTPException(400, "source_doc_id is required")

    # 1. Fetch target presentation to verify it exists
    target_doc = await get_document(id, user_id=user_id, db=db)
    if target_doc["doc_type"] != "presentation":
        raise HTTPException(400, "Target ID is not a presentation")

    # 2. Fetch source document (the note/text)
    source_doc = await get_document(source_id, user_id=user_id, db=db)
    source_text = source_doc.get("content_text") or ""
    
    if not source_text and source_doc.get("content"):
        # Fallback if it's a rich JSON structure (Tiptap)
        source_text = str(source_doc["content"]) # Simplistic, but better than nothing

    if not source_text:
        raise HTTPException(400, "Source document is empty")

    n_slides = body.get("n_slides", 10)

    prompt = f'''
# TASK: Transform the following document into a high-impact professional presentation.
Source Content:
---
{source_text[:6000]}
---

Target Number of Slides: {n_slides}
Theme: Wordex Sable & Cuivre (Atmospheric, Professional, Artistic)

# CONSTRAINTS:
1. Return ONLY pure valid JSON. 
2. NO conversational filler. NO markdown blocks.
3. Language: {source_doc.get('title', 'Same as source')}.

# STRUCTURE:
{{
  "title": "{source_doc['title']} - Pitch Deck",
  "theme": "sable",
  "slides": [
    {{
      "id": "s1",
      "title": "Slide Title",
      "content": "Key points from text...",
      "speakerNotes": "Context...",
      "layout": "titleAndContent"
    }}
  ]
}}
'''

    result = await ollama_generate(
        prompt=prompt,
        system=SYSTEM_PROMPTS.get("designer", ""),
        temperature=0.3,
    )
    
    try:
        # Clean potential markdown if LLM ignored instructions
        if "```" in result:
             result = result.split("```")[1].split("```")[0].strip()
             if result.startswith("json"): result = result[4:].strip()
             
        presentation_data = json.loads(result)
        for slide in presentation_data.get("slides", []):
            slide["id"] = str(uuid.uuid4())
                
    except Exception as e:
        raise HTTPException(500, f"AI generation failed to produce valid presentation JSON: {str(e)}")

    # 3. Save to target presentation
    update_data = DocumentUpdate(
        content=presentation_data, 
        title=presentation_data.get("title", f"Presentation: {source_doc['title']}")
    )
    await update_document(id, update_data, user_id=user_id, db=db)
    
    return {"status": "success", "presentation": presentation_data}

@router.get("/{id}/export-pptx")
async def export_pptx(
    id: str,
    user_id: str = Depends(get_current_user_id),
    db: asyncpg.Connection = Depends(get_db)
):
    try:
        from pptx import Presentation as PPTXPresentation
    except ImportError:
        raise HTTPException(501, "python-pptx is not installed on the server.")

    doc = await get_document(id, user_id=user_id, db=db)
    if doc["doc_type"] != "presentation":
        raise HTTPException(400, "Document is not a presentation")
    
    content = doc["content"] or {}
    slides_data = content.get("slides", [])
    title = content.get("title", doc["title"] or "Presentation")

    prs = PPTXPresentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title

    bullet_slide_layout = prs.slide_layouts[1]
    
    for sdata in slides_data:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        if title_shape:
            title_shape.text = sdata.get("title", "")
        
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = sdata.get("content", "")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    filename = f"{title.replace(' ', '_')}.pptx"
    
    return StreamingResponse(
        buf, 
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'}
    )
