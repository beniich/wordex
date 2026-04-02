import io
import json
from typing import Dict, Any

from docx import Document
from pptx import Presentation
from weasyprint import HTML
from jinja2 import Template

class ExportService:
    @staticmethod
    def _parse_tiptap_json(content: str) -> list:
        if not content:
            return []
        try:
            data = json.loads(content)
            return data.get("content", [])
        except Exception:
            return []

    @staticmethod
    def _extract_text(node: dict) -> str:
        text = ""
        for n in node.get("content", []):
            if n.get("type") == "text":
                text += n.get("text", "")
            elif n.get("type") == "paragraph":
                text += ExportService._extract_text(n)
        return text

    @staticmethod
    def generate_pdf(title: str, content_json: str, template: str = "default") -> bytes:
        nodes = ExportService._parse_tiptap_json(content_json)
        
        # HTML construction
        html_blocks = []
        for node in nodes:
            t = node.get("type")
            text = ExportService._extract_text(node)
            
            if t == "paragraph":
                html_blocks.append(f'<p class="body-text">{text}</p>')
            elif t == "heading":
                level = node.get("attrs", {}).get("level", 1)
                html_blocks.append(f"<h{level} class='heading-h{level}'>{text}</h{level}>")
            elif t == "bulletList":
                items = ""
                for li in node.get("content", []):
                    li_text = ExportService._extract_text(li)
                    items += f"<li>{li_text}</li>"
                html_blocks.append(f"<ul class='bullet-list'>{items}</ul>")
            elif t == "blockquote":
                html_blocks.append(f'<blockquote class="quote">{text}</blockquote>')

        body_html = "\n".join(html_blocks)

        template_str = """
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="utf-8">
            <title>{{ title }}</title>
            <link href="https://fonts.googleapis.com/css2?family=Manrope:wght@400;700;800&family=Playfair+Display:ital,wght@0,700;1,700&display=swap" rel="stylesheet">
            <style>
                @page { 
                    margin: 2.5cm; 
                    @bottom-right {
                        content: "Page " counter(page);
                        font-family: 'Manrope', sans-serif;
                        font-size: 8pt;
                        color: #a76526;
                    }
                    @bottom-left {
                        content: "Wordex Archives - {{ title }}";
                        font-family: 'Manrope', sans-serif;
                        font-size: 8pt;
                        text-transform: uppercase;
                        letter-spacing: 0.1em;
                        color: #a76526;
                    }
                }
                body { 
                    font-family: 'Manrope', sans-serif; 
                    color: #1c1c1a; 
                    background-color: #ffffff;
                    line-height: 1.7;
                }
                h1 { 
                    font-family: 'Playfair Display', serif;
                    font-size: 32pt;
                    color: #894d0d; 
                    margin-bottom: 30pt;
                    border-bottom: 3px solid #d8c3b4; 
                    padding-bottom: 10pt; 
                }
                .heading-h2 { 
                    font-family: 'Playfair Display', serif;
                    font-size: 22pt;
                    color: #524439; 
                    margin-top: 25pt;
                }
                .heading-h3 {
                    font-size: 16pt;
                    color: #a76526;
                    text-transform: uppercase;
                    letter-spacing: 0.05em;
                }
                .body-text { 
                    margin-bottom: 12pt;
                    font-size: 11pt;
                    text-align: justify;
                }
                .bullet-list { 
                    margin-left: 20pt; 
                    margin-bottom: 15pt;
                }
                li { margin-bottom: 6pt; }
                .quote {
                    font-style: italic;
                    border-left: 4px solid #d8c3b4;
                    padding-left: 15pt;
                    color: #524439;
                    margin: 20pt 0;
                    background: #faf7f2;
                    padding: 15pt;
                }
            </style>
        </head>
        <body>
            <h1>{{ title }}</h1>
            {{ body|safe }}
        </body>
        </html>
        """
        html_content = Template(template_str).render(title=title, body=body_html)
        
        # Generate PDF
        pdf_bytes = HTML(string=html_content).write_pdf()
        return pdf_bytes

    @staticmethod
    def generate_docx(title: str, content_json: str) -> bytes:
        nodes = ExportService._parse_tiptap_json(content_json)
        doc = Document()
        doc.add_heading(title, 0)

        for node in nodes:
            t = node.get("type")
            text = ExportService._extract_text(node)
            
            if t == "paragraph":
                doc.add_paragraph(text)
            elif t == "heading":
                level = node.get("attrs", {}).get("level", 1)
                doc.add_heading(text, level=level)
            elif t == "bulletList":
                for li in node.get("content", []):
                    li_text = ExportService._extract_text(li)
                    doc.add_paragraph(li_text, style='List Bullet')

        io_stream = io.BytesIO()
        doc.save(io_stream)
        io_stream.seek(0)
        return io_stream.read()

    @staticmethod
    def generate_pptx(title: str, content_json: str) -> bytes:
        nodes = ExportService._parse_tiptap_json(content_json)
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Content slides via basic markdown/H1 split
        bullet_slide_layout = prs.slide_layouts[1]
        
        current_slide = None
        current_tf = None
        
        for node in nodes:
            t = node.get("type")
            text = ExportService._extract_text(node)
            
            if t == "heading" and node.get("attrs", {}).get("level", 1) == 1:
                current_slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = current_slide.shapes
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                title_shape.text = text
                current_tf = body_shape.text_frame
            elif t == "heading" and current_tf:
                p = current_tf.add_paragraph()
                p.text = text
                p.level = 0
            elif t == "paragraph" and current_tf:
                p = current_tf.add_paragraph()
                p.text = text
                p.level = 0
            elif t == "bulletList" and current_tf:
                for li in node.get("content", []):
                    li_text = ExportService._extract_text(li)
                    p = current_tf.add_paragraph()
                    p.text = li_text
                    p.level = 1

        io_stream = io.BytesIO()
        prs.save(io_stream)
        io_stream.seek(0)
        return io_stream.read()
